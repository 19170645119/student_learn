"""PPTX 导出 — 将幻灯片JSON转为真正的PowerPoint文件（支持全部11种layout）"""
import json
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 配色主题
PRIMARY_COLOR = RGBColor(0x5B, 0x7F, 0xFF)
DARK_COLOR = RGBColor(0x33, 0x33, 0x33)
LIGHT_COLOR = RGBColor(0xF5, 0xF5, 0xF5)
WHITE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_COLOR = RGBColor(0xFF, 0x6B, 0x6B)
GREEN_COLOR = RGBColor(0x4A, 0xDE, 0x80)
PURPLE_COLOR = RGBColor(0xA7, 0x8B, 0xFA)

In = Inches

def _add_slide_with_title(prs, title_text, layout_index=1):
    """添加带标题的空白幻灯片"""
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.color.rgb = DARK_COLOR
    return slide

def _add_bullets(text_frame, bullets, font_size=16):
    """添加要点列表"""
    for i, bullet in enumerate(bullets):
        if i == 0 and text_frame.paragraphs[0].text:
            p = text_frame.paragraphs[0]
            p.clear()
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(8)
        p.level = 0

def _add_notes(slide, notes_text):
    """添加演讲者备注"""
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

def _add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_COLOR, bold=False, align=PP_ALIGN.LEFT):
    """添加文本框的便捷函数"""
    txBox = slide.shapes.add_textbox(In(left), In(top), In(width), In(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def _add_colored_bar(slide, top, height, color=PRIMARY_COLOR):
    """添加装饰色条"""
    shape = slide.shapes.add_shape(1, In(0), In(top), In(13.333), In(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def export_pptx(slides_data) -> bytes:
    """将幻灯片JSON数据转为PPTX文件字节流"""
    prs = Presentation()
    prs.slide_width = In(13.333)
    prs.slide_height = In(7.5)

    if isinstance(slides_data, str):
        slides_data = json.loads(slides_data)

    title = slides_data.get("title", "PPT课件")
    slides = slides_data.get("slides", [])

    for slide_data in slides:
        layout = slide_data.get("layout", "title_content")
        page_title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        table_data = slide_data.get("table")
        diagram_text = slide_data.get("diagram", "")
        key_point = slide_data.get("key_point", "")
        notes = slide_data.get("speaker_notes", "")
        section_num = slide_data.get("section_num", 0)
        section_title = slide_data.get("section_title", "")
        case_title = slide_data.get("case_title", "")
        case_content = slide_data.get("case_content", "")
        case_insight = slide_data.get("case_insight", "")
        metrics = slide_data.get("metrics", [])
        question = slide_data.get("question", "")
        hint = slide_data.get("hint", "")
        quote = slide_data.get("quote", "")
        author = slide_data.get("author", "")

        # ── 封面页 ──
        if layout == "title_slide":
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            _add_colored_bar(slide, 0, 0.15)
            _add_textbox(slide, 1.5, 2.0, 10, 2, title, font_size=44, color=PRIMARY_COLOR, bold=True, align=PP_ALIGN.CENTER)
            if bullets:
                _add_textbox(slide, 2, 4.2, 9, 1.5, " | ".join(bullets), font_size=20, color=DARK_COLOR, align=PP_ALIGN.CENTER)
            if key_point:
                _add_textbox(slide, 2, 5.5, 9, 1, "💬 " + key_point, font_size=16, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)
            _add_notes(slide, notes)
            continue

        # ── 章节过渡页 ──
        if layout == "chapter_divider":
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            # 背景色块
            bg_shape = slide.shapes.add_shape(1, In(1), In(1.5), In(11.333), In(4.5))
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            bg_shape.line.fill.background()
            _add_colored_bar(slide, 0, 0.15)
            if section_num:
                _add_textbox(slide, 2, 2.5, 9, 0.6, f"第{section_num}部分", font_size=16, color=PRIMARY_COLOR, bold=True, align=PP_ALIGN.CENTER)
            display_title = section_title or page_title
            _add_textbox(slide, 1.5, 3.2, 10, 1.5, display_title, font_size=36, color=RGBColor(0xE8,0xE8,0xF0), bold=True, align=PP_ALIGN.CENTER)
            if key_point:
                _add_textbox(slide, 2, 4.8, 9, 0.8, key_point, font_size=16, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)
            _add_notes(slide, notes)
            continue

        # ── 案例页 ──
        if layout == "case_study":
            slide = _add_slide_with_title(prs, case_title or page_title)
            if slide.placeholders:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1 and shape.has_text_frame:
                        tf = shape.text_frame
                        tf.clear()
                        # 案例标签
                        p = tf.paragraphs[0]
                        p.text = "📋 案例分析"
                        p.font.size = Pt(14)
                        p.font.color.rgb = ACCENT_COLOR
                        p.font.bold = True
                        # 案例内容
                        if case_content:
                            p2 = tf.add_paragraph()
                            p2.text = case_content
                            p2.font.size = Pt(16)
                            p2.font.color.rgb = DARK_COLOR
                            p2.space_after = Pt(12)
                        # 洞察
                        if case_insight:
                            p3 = tf.add_paragraph()
                            p3.text = f"💡 洞察：{case_insight}"
                            p3.font.size = Pt(14)
                            p3.font.color.rgb = GREEN_COLOR
                            p3.font.bold = True
                        # 补充要点
                        if bullets:
                            for b in bullets:
                                pb = tf.add_paragraph()
                                pb.text = f"• {b}"
                                pb.font.size = Pt(14)
                        break
            _add_notes(slide, notes)
            continue

        # ── 数据亮点页 ──
        if layout == "data_highlight":
            slide = _add_slide_with_title(prs, page_title)
            if metrics and slide.placeholders:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1 and shape.has_text_frame:
                        tf = shape.text_frame
                        tf.clear()
                        for m in metrics:
                            p = tf.add_paragraph()
                            p.text = f"{m.get('value', '')}  —  {m.get('label', '')}"
                            p.font.size = Pt(22)
                            p.font.color.rgb = PRIMARY_COLOR
                            p.font.bold = True
                            p.space_after = Pt(4)
                            if m.get('desc'):
                                pd = tf.add_paragraph()
                                pd.text = m['desc']
                                pd.font.size = Pt(14)
                                pd.font.color.rgb = DARK_COLOR
                                pd.space_after = Pt(12)
                        if bullets:
                            for b in bullets:
                                pb = tf.add_paragraph()
                                pb.text = f"• {b}"
                                pb.font.size = Pt(14)
                        break
            _add_notes(slide, notes)
            continue

        # ── 互动提问页 ──
        if layout == "audience_question":
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            _add_textbox(slide, 2, 1.5, 9, 0.8, "❓ 互动提问", font_size=18, color=ACCENT_COLOR, bold=True, align=PP_ALIGN.CENTER)
            _add_textbox(slide, 1.5, 2.8, 10, 1.5, question or page_title, font_size=28, color=DARK_COLOR, bold=True, align=PP_ALIGN.CENTER)
            if hint:
                _add_textbox(slide, 2, 4.5, 9, 1, f"💡 提示：{hint}", font_size=16, color=GREEN_COLOR, align=PP_ALIGN.CENTER)
            if key_point:
                _add_textbox(slide, 2, 5.5, 9, 0.6, key_point, font_size=14, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.CENTER)
            _add_notes(slide, notes)
            continue

        # ── 金句引用页 ──
        if layout == "key_quote":
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            _add_textbox(slide, 2, 1.0, 9, 0.8, "\u201C", font_size=72, color=PRIMARY_COLOR, align=PP_ALIGN.CENTER)
            _add_textbox(slide, 1.5, 2.2, 10, 2, quote or page_title, font_size=26, color=DARK_COLOR, bold=True, align=PP_ALIGN.CENTER)
            if author:
                _add_textbox(slide, 2, 4.5, 9, 0.6, f"—— {author}", font_size=16, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)
            if key_point:
                _add_textbox(slide, 2, 5.3, 9, 0.6, key_point, font_size=14, color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.CENTER)
            _add_notes(slide, notes)
            continue

        # ── 普通内容页（title_content / two_column / summary / diagram / table）──
        slide = _add_slide_with_title(prs, page_title)

        if slide.placeholders:
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break

            if body_shape and body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()

                if layout == "title_content":
                    _add_bullets(tf, bullets)
                    if key_point:
                        p = tf.add_paragraph()
                        p.text = f"★ {key_point}"
                        p.font.size = Pt(14)
                        p.font.color.rgb = PRIMARY_COLOR
                        p.font.bold = True

                elif layout == "two_column":
                    mid = (len(bullets) + 1) // 2
                    p = tf.paragraphs[0]
                    p.text = "  |  ".join(bullets[:mid])
                    p.font.size = Pt(16)
                    p2 = tf.add_paragraph()
                    p2.text = "  |  ".join(bullets[mid:])
                    p2.font.size = Pt(16)

                elif layout == "summary":
                    _add_bullets(tf, bullets, font_size=18)
                    if tf.paragraphs:
                        tf.paragraphs[0].font.bold = True

                elif layout == "diagram":
                    if diagram_text:
                        p = tf.paragraphs[0]
                        p.text = diagram_text
                        p.font.size = Pt(14)
                        p.font.color.rgb = DARK_COLOR
                    if bullets:
                        for b in bullets:
                            p = tf.add_paragraph()
                            p.text = f"• {b}"
                            p.font.size = Pt(14)

        # 表格（独立处理）
        if layout == "table" and table_data:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            if headers and rows:
                num_rows = len(rows) + 1
                num_cols = len(headers)
                tbl = slide.shapes.add_table(num_rows, num_cols, In(1), In(2.5), In(11), In(0.5 * num_rows)).table
                for col_idx, header in enumerate(headers):
                    cell = tbl.cell(0, col_idx)
                    cell.text = header
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(14)
                        p.font.bold = True
                        p.font.color.rgb = WHITE_COLOR
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PRIMARY_COLOR
                for row_idx, row in enumerate(rows):
                    for col_idx, val in enumerate(row):
                        cell = tbl.cell(row_idx + 1, col_idx)
                        cell.text = str(val)
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(13)

        _add_notes(slide, notes)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()
